import asyncio
import dataclasses
import logging
import re
from pathlib import Path
from typing import Any

from .config import config
from .db import connect, dumps, load_llm_settings
from .llm import embed_texts, summarize_source
from .vector_store import delete_source as delete_source_vectors
from .vector_store import upsert_chunks


ALLOWED_EXTENSIONS = {
    ".pdf", ".txt", ".md", ".markdown", ".docx", ".html", ".htm", ".srt", ".vtt",
    ".xlsx", ".csv", ".pptx",
}
logger = logging.getLogger(__name__)

# Inline WebVTT tags (e.g. <v Speaker>, <00:00:01.000>, <c.classname>) stripped
# from caption text so only the spoken words remain.
_VTT_INLINE_TAG = re.compile(r"<[^>]+>")


def supported(filename: str) -> bool:
    """Return whether the filename extension is accepted for ingestion."""
    return Path(filename).suffix.lower() in ALLOWED_EXTENSIONS


@dataclasses.dataclass
class ExtractionResult:
    """Extracted sections plus what the extractor noticed on the way (A6a).

    ``extractor`` and ``notes`` capture things only the extractor itself can
    know — most importantly *which* code path produced the sections, since a
    PDF that fell back to plain pypdf loses tables and paragraph structure and
    therefore yields coarser citations. Diagnostics that can be measured from
    the sections afterwards (counts, char totals) are **not** carried here;
    they are derived in `collect_ingest_diagnostics` so there is one place that
    computes them.
    """
    sections: list[tuple[str, str]]
    extractor: str
    notes: list[str] = dataclasses.field(default_factory=list)
    #: Format-specific facts for the diagnostics panel (sheet types, header
    #: decision, CSV encoding …). Free-form so a new format adds signals without
    #: a schema change; rendered generically.
    details: dict[str, Any] = dataclasses.field(default_factory=dict)
    #: True when ``sections`` are already chunk-shaped and must NOT be re-packed
    #: by :func:`chunk_sections`. Spreadsheets set this: a row is the semantic
    #: unit, and sentence-packing across rows would glue unrelated records
    #: together and destroy the ``sheet "X" row N`` citation labels.
    pre_chunked: bool = False


def extract_sections(path: Path) -> ExtractionResult:
    """Extract text sections from a supported source file.

    Dispatches per suffix to a helper. Each helper returns a list of
    ``(location, text)`` pairs; the location label flows through to chunk
    citations so users can see whether an answer came from the body, a
    header, a footnote, etc.
    """
    suffix = path.suffix.lower()
    logger.info("extract_started path=%s suffix=%s", path.name, suffix)
    notes: list[str] = []
    if suffix == ".pdf":
        # Structured first; note the degradation explicitly when falling back,
        # because "why are this file's citations only page-level?" is otherwise
        # invisible to the user.
        sections = _extract_pdf_with_pdfplumber(path)
        if sections:
            extractor = "pdf_pdfplumber"
        else:
            sections = _extract_pdf_with_pypdf(path)
            extractor = "pdf_pypdf"
            notes.append("pdf_structure_fallback")
    elif suffix == ".docx":
        sections = _extract_docx(path)
        extractor = "docx"
    elif suffix in {".html", ".htm"}:
        sections = _extract_html(path)
        extractor = "html"
    elif suffix in {".srt", ".vtt"}:
        sections = _extract_subtitles(path)
        extractor = "subtitles"
    elif suffix in {".xlsx", ".csv"}:
        # Spreadsheets return fully-formed chunks (see _extract_spreadsheet).
        return _extract_spreadsheet(path)
    elif suffix == ".pptx":
        # Slide sections flow through the normal chunker (short slides pack
        # together rather than becoming one tiny vector each).
        return _extract_pptx(path)
    else:
        sections = [("document", path.read_text(encoding="utf-8", errors="ignore"))]
        extractor = "plain_text"
    logger.info(
        "extract_completed path=%s sections=%s extractor=%s", path.name, len(sections), extractor
    )
    return ExtractionResult(sections=sections, extractor=extractor, notes=notes)


def _extract_subtitles(path: Path) -> list[tuple[str, str]]:
    """Extract spoken text from an .srt / .vtt subtitle file (A7).

    Strips cue index numbers, timestamp lines, the WebVTT header and
    NOTE/STYLE/REGION metadata blocks, and inline VTT tags — leaving the
    caption text as a single ``transcript`` section. Consecutive duplicate
    lines (common with rolling captions) are collapsed. No new dependency.
    """
    raw = path.read_text(encoding="utf-8", errors="ignore")
    lines: list[str] = []
    skip_block = False
    for line in raw.splitlines():
        stripped = line.strip()
        if not stripped:
            skip_block = False  # a blank line ends any NOTE/STYLE block
            continue
        # WebVTT header + metadata blocks run until the next blank line.
        if stripped == "WEBVTT" or stripped.startswith(("WEBVTT", "NOTE", "STYLE", "REGION")):
            skip_block = True
            continue
        if skip_block:
            continue
        if "-->" in stripped:  # timestamp cue line (SRT or VTT, incl. positioning)
            continue
        if stripped.isdigit():  # bare SRT cue index
            continue
        text = _VTT_INLINE_TAG.sub("", stripped).strip()
        if not text:
            continue
        if lines and lines[-1] == text:  # collapse rolling-caption repeats
            continue
        lines.append(text)
    transcript = "\n".join(lines)
    return [("transcript", transcript)] if transcript else []


def _extract_pdf_with_pypdf(path: Path) -> list[tuple[str, str]]:
    """Plain page-level PDF text — the fallback when pdfplumber can't help.

    Loses tables and paragraph blocks, so citations degrade to ``page N``.
    `extract_sections` records this as the ``pdf_structure_fallback`` note.
    """
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return [
        (f"page {index}", page.extract_text() or "")
        for index, page in enumerate(reader.pages, start=1)
    ]


def _extract_pdf_with_pdfplumber(path: Path) -> list[tuple[str, str]]:
    """Best-effort structured PDF extraction using pdfplumber."""
    try:
        import pdfplumber
    except Exception:
        return []

    try:
        with pdfplumber.open(str(path)) as pdf:
            sections: list[tuple[str, str]] = []
            for page_index, page in enumerate(pdf.pages, start=1):
                sections.extend(_extract_pdf_page_blocks(page, page_index))
            return sections
    except Exception:
        logger.exception("pdfplumber_extract_failed path=%s", path.name)
        return []


def _extract_pdf_page_blocks(page, page_index: int) -> list[tuple[str, str]]:
    """Extract table + paragraph blocks from a page, preserving top-to-bottom order."""
    tables = list(page.find_tables() or [])
    table_bboxes = [table.bbox for table in tables if getattr(table, "bbox", None)]

    table_blocks: list[dict[str, Any]] = []
    for table_index, table in enumerate(tables, start=1):
        table_text = _render_pdf_table(table.extract() or [])
        if not table_text:
            continue
        table_blocks.append(
            {
                "top": _pdf_bbox_top(getattr(table, "bbox", None)),
                "location": f"page {page_index} table {table_index}",
                "text": table_text,
            }
        )

    words = list(page.extract_words() or [])
    non_table_words = [
        word for word in words if not _pdf_word_in_any_bbox(word, table_bboxes)
    ]
    paragraph_blocks = _pdf_words_to_paragraph_blocks(non_table_words, page_index)

    merged = [
        {"kind": "paragraph", **block}
        for block in paragraph_blocks
    ] + [
        {"kind": "table", **block}
        for block in table_blocks
    ]
    merged.sort(key=lambda block: (block["top"], 0 if block["kind"] == "paragraph" else 1))

    return [
        (block["location"], block["text"])
        for block in merged
        if block["text"].strip()
    ]


def _pdf_bbox_top(bbox: tuple[float, float, float, float] | None) -> float:
    if not bbox:
        return 0.0
    return float(bbox[1])


def _pdf_word_in_any_bbox(
    word: dict[str, Any],
    bboxes: list[tuple[float, float, float, float]],
) -> bool:
    x0 = float(word.get("x0", 0.0))
    x1 = float(word.get("x1", x0))
    top = float(word.get("top", 0.0))
    bottom = float(word.get("bottom", top))
    cx = (x0 + x1) / 2.0
    cy = (top + bottom) / 2.0
    for bx0, btop, bx1, bbottom in bboxes:
        if bx0 <= cx <= bx1 and btop <= cy <= bbottom:
            return True
    return False


def _pdf_words_to_paragraph_blocks(words: list[dict[str, Any]], page_index: int) -> list[dict[str, Any]]:
    """Group non-table words into paragraph blocks by line and vertical gap."""
    if not words:
        return []

    sorted_words = sorted(words, key=lambda w: (float(w.get("top", 0.0)), float(w.get("x0", 0.0))))
    line_tolerance = 3.0
    lines: list[dict[str, Any]] = []
    current_line: list[dict[str, Any]] = []
    line_top = 0.0

    def flush_line() -> None:
        nonlocal current_line, line_top
        if not current_line:
            return
        current_line.sort(key=lambda w: float(w.get("x0", 0.0)))
        text = " ".join((w.get("text") or "").strip() for w in current_line if (w.get("text") or "").strip())
        if text:
            tops = [float(w.get("top", line_top)) for w in current_line]
            bottoms = [float(w.get("bottom", line_top)) for w in current_line]
            lines.append(
                {
                    "top": min(tops),
                    "bottom": max(bottoms),
                    "text": text,
                }
            )
        current_line = []

    for word in sorted_words:
        word_top = float(word.get("top", 0.0))
        if not current_line:
            current_line = [word]
            line_top = word_top
            continue
        if abs(word_top - line_top) <= line_tolerance:
            current_line.append(word)
        else:
            flush_line()
            current_line = [word]
            line_top = word_top
    flush_line()

    paragraph_gap = 12.0
    paragraphs: list[dict[str, Any]] = []
    para_lines: list[str] = []
    para_top = 0.0
    previous_bottom = 0.0
    paragraph_index = 1

    def flush_paragraph() -> None:
        nonlocal para_lines, para_top, paragraph_index
        if not para_lines:
            return
        text = "\n".join(para_lines).strip()
        if text:
            paragraphs.append(
                {
                    "top": para_top,
                    "location": f"page {page_index} paragraph {paragraph_index}",
                    "text": text,
                }
            )
            paragraph_index += 1
        para_lines = []

    for line in lines:
        if not para_lines:
            para_top = float(line["top"])
            para_lines = [line["text"]]
            previous_bottom = float(line["bottom"])
            continue
        gap = float(line["top"]) - previous_bottom
        if gap > paragraph_gap:
            flush_paragraph()
            para_top = float(line["top"])
            para_lines = [line["text"]]
        else:
            para_lines.append(line["text"])
        previous_bottom = float(line["bottom"])
    flush_paragraph()

    return paragraphs


def _render_pdf_table(rows: list[list[Any]]) -> str:
    """Render extracted table rows into retrieval-friendly pipe-separated text."""
    cleaned_rows: list[list[str]] = []
    for row in rows:
        cleaned = [" ".join(str(cell or "").split()) for cell in row]
        if any(cell for cell in cleaned):
            cleaned_rows.append(cleaned)
    if not cleaned_rows:
        return ""
    return "Table:\n" + "\n".join(" | ".join(row) for row in cleaned_rows)


def _extract_docx(path: Path) -> list[tuple[str, str]]:
    """Extract paragraphs, tables, headers, footers, text boxes, footnotes.

    The historical implementation only walked ``doc.paragraphs``, silently
    dropping every cell in every table (which in a typical case-study docx
    is 99% of the content). This walks the document body in order and
    flattens tables inline (with nested-table recursion), then emits
    headers / footers / text boxes / footnotes as their own labelled
    sections so citations can tell users where evidence came from.
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(str(path))
    sections: list[tuple[str, str]] = []

    body_text = _render_docx_container(doc)
    if body_text.strip():
        sections.append(("document", body_text))

    # Headers / footers are in their own XML parts, not in doc.element.body.
    header_chunks, footer_chunks = [], []
    for section in doc.sections:
        h = _render_docx_container(section.header)
        f = _render_docx_container(section.footer)
        if h.strip():
            header_chunks.append(h)
        if f.strip():
            footer_chunks.append(f)
    if header_chunks:
        sections.append(("header", "\n\n".join(header_chunks)))
    if footer_chunks:
        sections.append(("footer", "\n\n".join(footer_chunks)))

    # Text boxes (<w:txbxContent>) live inside drawings and are skipped by
    # the body's CT_P / CT_Tbl iteration. Flatten any w:t runs we find.
    txbx_chunks = []
    for txbx in doc.element.iter(qn("w:txbxContent")):
        text = " ".join(t.text for t in txbx.iter(qn("w:t")) if t.text)
        if text.strip():
            txbx_chunks.append(text)
    if txbx_chunks:
        sections.append(("text boxes", "\n\n".join(txbx_chunks)))

    # Footnotes / endnotes are stored as separate package parts referenced
    # from the document part by relationship type. Each w:t inside the part
    # is a footnote body run; we just concatenate.
    note_chunks = []
    for rel in doc.part.rels.values():
        if "footnote" in rel.reltype or "endnote" in rel.reltype:
            try:
                root = rel.target_part.element
            except AttributeError:
                continue
            text = " ".join(t.text for t in root.iter(qn("w:t")) if t.text)
            if text.strip():
                note_chunks.append(text)
    if note_chunks:
        sections.append(("footnotes", "\n\n".join(note_chunks)))

    return sections


def _iter_docx_block_items(parent):
    """Yield Paragraph and Table children of parent in document order.

    Works for Document (the body), _Cell (cell contents, for nested tables),
    and _Header / _Footer (their own XML root). Order matters: python-docx's
    ``parent.paragraphs`` and ``parent.tables`` each return a flat list, so
    a doc that alternates paragraphs and tables would lose its narrative
    flow if you combined them naively.
    """
    from docx.document import Document as _Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph

    if isinstance(parent, _Document):
        parent_elem = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elem = parent._tc
    elif hasattr(parent, "_element"):
        # _Header / _Footer
        parent_elem = parent._element
    else:
        parent_elem = parent

    for child in parent_elem.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _render_docx_container(container) -> str:
    """Render a Document / _Header / _Footer / _Cell as text with tables inline."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    parts: list[str] = []
    for block in _iter_docx_block_items(container):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if text:
                parts.append(text)
        elif isinstance(block, Table):
            rendered = _render_docx_table(block)
            if rendered:
                parts.append(rendered)
    return "\n".join(parts)


def _render_docx_table(table) -> str:
    """Render a table as ``Table:`` plus ' | '-separated cells per row.

    Nested tables (a table inside a cell) are rendered inline by recursing
    through ``_render_docx_container`` on each cell. Whitespace inside a
    cell is collapsed to single spaces so the row separator stays
    unambiguous.
    """
    rows: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            cell_text = _render_docx_container(cell)
            cells.append(" ".join(cell_text.split()))
        if any(c.strip() for c in cells):
            rows.append(" | ".join(cells))
    if not rows:
        return ""
    return "Table:\n" + "\n".join(rows)


def _extract_html(path: Path) -> list[tuple[str, str]]:
    """Strip noise + recover alt / title / meta-description text BeautifulSoup
    would otherwise drop. Returns a single section.

    What changed vs the naive ``soup.get_text``:
      - script / style / noscript / template removed (noise).
      - elements with ``hidden`` attribute or inline ``style='display:none'``
        removed (catches injected honeypots / hidden JSON-LD blocks).
      - <meta name='description'> and og:description appended (cheap recall).
      - <img alt> / <a title> / <input value> appended as ``[image: ...]``
        style sidecar lines so they survive get_text() without polluting
        the main flow.
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    for element in soup(["script", "style", "noscript", "template"]):
        element.decompose()
    # Hidden via attribute or inline style. We deliberately do NOT touch
    # aria-hidden, visibility:hidden, or CSS-class-based hiding — those
    # often mark legitimate collapsed content (tabs, accordions) that the
    # user can still expand to read.
    for element in soup.find_all(hidden=True):
        element.decompose()
    for element in soup.find_all(style=True):
        style = element.get("style", "").lower().replace(" ", "")
        if "display:none" in style:
            element.decompose()

    extras: list[str] = []
    meta_desc = soup.find("meta", attrs={"name": "description"})
    if meta_desc and meta_desc.get("content"):
        extras.append(meta_desc["content"].strip())
    og_desc = soup.find("meta", attrs={"property": "og:description"})
    if og_desc and og_desc.get("content"):
        extras.append(og_desc["content"].strip())
    for img in soup.find_all("img", alt=True):
        alt = (img.get("alt") or "").strip()
        if alt:
            extras.append(f"[image: {alt}]")
    for anchor in soup.find_all("a", title=True):
        title = (anchor.get("title") or "").strip()
        if title:
            extras.append(f"[link: {title}]")
    for inp in soup.find_all("input", value=True):
        if inp.get("type") in {"hidden", "password", "submit", "button"}:
            continue
        value = (inp.get("value") or "").strip()
        if value:
            extras.append(f"[input: {value}]")

    body_text = soup.get_text("\n")
    if extras:
        body_text = body_text + "\n\n" + "\n".join(extras)
    return [("html", body_text)]


# Sentence boundary regex. Matches the END position after a terminator:
#   - CJK terminators (。！？): no trailing-space requirement, since CJK
#     doesn't space-separate sentences.
#   - Latin terminators (.!?): must be followed by whitespace or end of
#     input, to avoid splitting decimals (3.14), URLs, and most abbreviations.
#   - One or more newlines: always a boundary.
# Known limitation: "Mr. Smith" still splits at "Mr.". Acceptable for POC.
_SENTENCE_BOUNDARY_RE = re.compile(r"[。！？]+|[.!?](?=\s|$)|\n+")
# Soft punctuation used as a fallback when a single sentence is longer than
# the target chunk size. Includes both CJK and Latin commas / semicolons.
_SOFT_BREAK_RE = re.compile(r"[，、；,;]")
_CJK_RE = re.compile(r"[一-鿿]")
# Internal whitespace normalisation: collapse runs but preserve newlines so
# split_sentences() can use them as boundaries. PDFs often inject erratic
# spacing inside paragraphs that we still want flattened.
_HORIZONTAL_WS_RE = re.compile(r"[ \t\r\f\v]+")

# Chunking targets (config-driven; changing them requires re-indexing).
LATIN_TARGET_CHARS = config.chunking.latin_target_chars
CJK_TARGET_CHARS = config.chunking.cjk_target_chars
DEFAULT_OVERLAP_SENTENCES = config.chunking.overlap_sentences


#: Formats whose on-disk size does not bound their parse cost — each for its own
#: reason, so do not collapse these into one sentence:
#:
#: * ``.xlsx`` — zip expansion only. openpyxl opens it ``read_only=True``, so the
#:   rows themselves stream; the risk is a small archive decompressing to
#:   gigabytes.
#: * ``.pptx`` — zip expansion, *plus* python-pptx materialises the whole
#:   presentation object rather than streaming slides.
#: * ``.csv`` — not an archive and now parsed incrementally, but the cap still
#:   bounds parser CPU and pathological single records/fields that can approach
#:   the file's full size even when the surrounding rows stream.
#:
#: PDF/DOCX stream and are not listed.
STRICT_EXTRACT_CAP_SUFFIXES = frozenset({".xlsx", ".pptx", ".csv"})


def upload_limit_for(filename: str) -> int:
    """Effective per-file upload cap for this filename, in bytes.

    Normally `upload_max_file_bytes`, but formats with the stricter parser cap
    fall back to `extract_max_file_bytes` because that is what the worker would
    enforce anyway. Applying it at upload time turns a confusing two-stage
    rejection — upload succeeds, then ingest fails minutes later — into one
    immediate, explainable error.

    Lives here rather than in `app/main.py` so the "which formats are expensive
    to parse" knowledge stays next to the parsers that make it true.
    """
    limit = config.runtime.upload_max_file_bytes
    if Path(filename).suffix.lower() in STRICT_EXTRACT_CAP_SUFFIXES:
        return min(limit, config.runtime.extract_max_file_bytes)
    return limit


def _guard_source_size(path: Path) -> None:
    """Refuse sources whose parser cost needs a stricter size backstop.

    The upload path already applies the same cap (`upload_limit_for`), so in
    normal operation this never fires. It stays as the backstop for files that
    did not arrive through an upload — reindexing a source stored before the cap
    existed, or a file placed in `data/uploads/` by hand — because this is the
    check that actually stands between archive expansion / oversized records
    and their parsers.
    """
    limit = config.runtime.extract_max_file_bytes
    size = path.stat().st_size
    if size > limit:
        raise ValueError(f"Source file is too large to ingest ({size} bytes > {limit}).")


def _iter_pptx_shapes(shapes):
    """Yield shapes depth-first, descending into groups.

    Grouped shapes are the PPTX equivalent of the nested-table bug we hit in
    DOCX: their text is invisible to a flat `slide.shapes` walk, so a deck whose
    author grouped its bullet boxes would index as an empty slide.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def _extract_pptx(path: Path) -> ExtractionResult:
    """Extract a .pptx deck as slide-scoped sections (A6b Phase 1, text-first).

    Emits `slide N` (title + body), `slide N table K`, and `slide N notes` so a
    citation points at a specific slide. Images/diagrams are **not** read — that
    needs OCR (`A8`) or a vision model (`A9`); slides that carry only visual
    content are reported through diagnostics instead of silently indexing as
    empty, which is the whole point of shipping Phase 1 before either exists.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    _guard_source_size(path)
    presentation = Presentation(str(path))
    sections: list[tuple[str, str]] = []
    details: dict[str, Any] = {
        "slides": 0,
        "slides_without_text": 0,
        "tables": 0,
        "notes": 0,
        "images": 0,
    }
    notes_list: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        details["slides"] += 1
        title_shape = slide.shapes.title
        lines: list[str] = []
        table_sections: list[tuple[str, str]] = []
        table_number = 0
        has_visual = False
        for shape in _iter_pptx_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                table_number += 1
                details["tables"] += 1
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                rendered = _render_pdf_table(rows)
                if rendered:
                    # Held back so each slide reads body -> tables -> notes;
                    # PPTX shape order does not follow visual reading order.
                    table_sections.append((f"slide {index} table {table_number}", rendered))
                continue
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                details["images"] += 1
                has_visual = True
                continue
            if not getattr(shape, "has_text_frame", False):
                # Charts, media, SmartArt — visual content Phase 1 cannot read.
                has_visual = True
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title_shape is not None and shape is title_shape:
                lines.insert(0, text)
            else:
                lines.append(text)

        body = "\n".join(lines).strip()
        if body:
            sections.append((f"slide {index}", body))
        sections.extend(table_sections)
        if not body and not table_sections and has_visual:
            # A slide with visuals but no readable text: the exact case a user
            # would otherwise blame on retrieval.
            details["slides_without_text"] += 1

        if getattr(slide, "has_notes_slide", False):
            note_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if note_text:
                details["notes"] += 1
                sections.append((f"slide {index} notes", note_text))

    if details["slides_without_text"]:
        notes_list.append("pptx_visual_only_slides")
    logger.info(
        "extract_completed path=%s sections=%s extractor=pptx slides=%s",
        path.name, len(sections), details["slides"],
    )
    return ExtractionResult(
        sections=sections, extractor="pptx", notes=notes_list, details=details
    )


# --------------------------------------------------------------------------
# A6c · Spreadsheet ingestion (.xlsx / .csv)
#
# Design decisions (why, not just what) live in docs/SPREADSHEET_INGESTION.md.
# The MVP detects Q&A sheets only; every other shape falls back to bounded
# generic-record chunking and *says so* in diagnostics rather than pretending
# it understood the sheet. All workbook access goes through _read_xlsx_sheets /
# _read_csv_sheet so the reader can be swapped (python-calamine) without
# touching chunk shaping.
# --------------------------------------------------------------------------

def _column_label(index: int) -> str:
    """Excel-style fallback name for a column with no header (0 -> 'Column A')."""
    letters = ""
    n = index
    while True:
        letters = chr(ord("A") + n % 26) + letters
        n = n // 26 - 1
        if n < 0:
            break
    return f"Column {letters}"


def _cell_text(value: Any) -> str:
    """Render a cell as trimmed display text ('' for blanks)."""
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _looks_like_header(rows: list[list[Any]]) -> bool:
    """Decide whether row 0 is a header rather than data.

    Deliberately conservative — mislabelling data as a header silently deletes
    a record, so every signal must point the same way: all cells textual, no
    duplicates, and none long enough to be prose.
    """
    if not rows:
        return False
    values = [v for v in rows[0] if _cell_text(v)]
    if not values:
        return False
    if any(isinstance(v, (int, float)) and not isinstance(v, bool) for v in values):
        return False
    labels = [_cell_text(v) for v in values]
    if len(set(labels)) != len(labels):
        return False
    if any(len(label) > 40 for label in labels):
        return False
    # Sentences are data, not column names — this is what separates a headerless
    # exported FAQ ("忘記密碼怎麼辦？") from a real header row ("問題").
    if any(label[-1] in "。？！?!." for label in labels):
        return False
    # A header should also be terser than the rows beneath it.
    body = [_cell_text(v) for row in rows[1:] for v in row if _cell_text(v)]
    if body:
        avg_header = sum(len(label) for label in labels) / len(labels)
        avg_body = sum(len(text) for text in body) / len(body)
        if avg_body and avg_header > avg_body * 1.2:
            return False
    return True


def _synonyms(raw: str) -> list[str]:
    return [part.strip().lower() for part in (raw or "").split(",") if part.strip()]


def _match_column(columns: list[str], synonyms: list[str]) -> int | None:
    """Find the column whose name matches one of `synonyms`.

    Short synonyms ('q', 'a') must match exactly — as a substring they would hit
    almost any English header. Longer ones may match as a substring so house
    vocabulary like 客戶提問 still resolves via 提問.
    """
    normalized = [c.strip().lower() for c in columns]
    for synonym in synonyms:
        for index, name in enumerate(normalized):
            if name == synonym:
                return index
    for synonym in synonyms:
        if len(synonym) <= 2:
            continue
        for index, name in enumerate(normalized):
            if synonym in name:
                return index
    return None


def _detect_qa_columns(columns: list[str], has_header: bool) -> dict[str, Any] | None:
    """Return the Q&A column mapping for this sheet, or None if it isn't Q&A."""
    s = config.spreadsheet
    question = _match_column(columns, _synonyms(s.qa_question_synonyms))
    answer = _match_column(columns, _synonyms(s.qa_answer_synonyms))
    if question is not None and answer is not None and question != answer:
        return {"question": question, "answer": answer, "auto_detected": False}
    # A headerless two-column sheet is the classic exported FAQ: treat it as
    # question/answer but record that the mapping was guessed.
    if not has_header and len(columns) == 2:
        return {"question": 0, "answer": 1, "auto_detected": True}
    return None


def _meta_columns(columns: list[str], qa: dict[str, Any]) -> list[int]:
    """Non-Q&A columns worth embedding as context (category / tags / keywords)."""
    wanted = ("category", "分類", "類別", "tag", "標籤", "keyword", "關鍵字")
    out = []
    for index, name in enumerate(columns):
        if index in (qa["question"], qa["answer"]):
            continue
        lowered = name.strip().lower()
        if any(token in lowered for token in wanted):
            out.append(index)
    return out


def _qa_sections(
    sheet: str, columns: list[str], rows: list[tuple[int, list[Any]]], qa: dict[str, Any]
) -> list[tuple[str, str]]:
    """One chunk per Q&A row, in the trimmed-preamble shape (design doc)."""
    meta = _meta_columns(columns, qa)
    sections: list[tuple[str, str]] = []
    for row_number, values in rows:
        question = _cell_text(values[qa["question"]]) if qa["question"] < len(values) else ""
        answer = _cell_text(values[qa["answer"]]) if qa["answer"] < len(values) else ""
        if not question and not answer:
            continue
        preamble = f"Sheet: {sheet}"
        for index in meta:
            text = _cell_text(values[index]) if index < len(values) else ""
            if text:
                preamble += f" · {columns[index]}: {text}"
        body = f"{preamble}\n\nQuestion:\n{question}\n\nAnswer:\n{answer}"
        sections.append((f'sheet "{sheet}" row {row_number}', body))
    return sections


def _record_row_text(row_number: int, columns: list[str], values: list[Any]) -> str:
    """`Row N:` block with one `column = value` line per populated cell."""
    lines = [f"Row {row_number}:"]
    for index, name in enumerate(columns):
        text = _cell_text(values[index]) if index < len(values) else ""
        if text:
            lines.append(f"{name} = {text}")
    return "\n".join(lines)


def _split_wide_row(
    sheet: str, columns: list[str], values: list[Any], row_number: int, preamble: str, budget: int
) -> list[tuple[str, str]]:
    """Split one over-budget row into column groups that repeat the identifier.

    Without this a very wide row would be embedded and silently truncated past
    the model's window, leaving its tail columns unreachable by vector search.
    Every part repeats column 0 so each child chunk still identifies its record.
    """
    identifier = ""
    if columns:
        head = _cell_text(values[0]) if values else ""
        identifier = f"{columns[0]} = {head}" if head else ""
    groups: list[list[str]] = []
    current: list[str] = []
    for index, name in enumerate(columns[1:], start=1):
        text = _cell_text(values[index]) if index < len(values) else ""
        if not text:
            continue
        line = f"{name} = {text}"
        candidate = current + [line]
        probe = f"{preamble}\n\nRow {row_number}:\n{identifier}\n" + "\n".join(candidate)
        if current and estimate_embedding_tokens(probe) > budget:
            groups.append(current)
            current = [line]
        else:
            current = candidate
    if current:
        groups.append(current)
    if not groups:
        groups = [[]]

    total = len(groups)
    sections: list[tuple[str, str]] = []
    for part, lines in enumerate(groups, start=1):
        header = [f"Row {row_number}:"]
        if identifier:
            header.append(identifier)
        body = f"{preamble}\n\n" + "\n".join(header + lines)
        location = f'sheet "{sheet}" row {row_number}'
        if total > 1:
            location += f" part {part}/{total}"
        sections.append((location, body))
    return sections


def _record_sections(
    sheet: str, columns: list[str], rows: list[tuple[int, list[Any]]]
) -> list[tuple[str, str]]:
    """Generic-record chunking with token-aware adaptive row packing.

    Rows are packed until the *estimated* embedding tokens reach the budget, so
    wide sheets naturally degrade to one row per chunk instead of producing
    chunks whose tails would be truncated away.
    """
    s = config.spreadsheet
    budget = s.embed_token_budget
    preamble = f"Sheet: {sheet} · Columns: {', '.join(columns)}"
    sections: list[tuple[str, str]] = []
    buffer: list[tuple[int, str]] = []

    def flush() -> None:
        if not buffer:
            return
        first, last = buffer[0][0], buffer[-1][0]
        location = (
            f'sheet "{sheet}" row {first}' if first == last
            else f'sheet "{sheet}" rows {first}-{last}'
        )
        body = preamble + "\n\n" + "\n\n".join(text for _, text in buffer)
        sections.append((location, body))
        buffer.clear()

    for row_number, values in rows:
        row_text = _record_row_text(row_number, columns, values)
        if row_text.count("\n") == 0:  # header line only -> the row was blank
            continue
        if estimate_embedding_tokens(f"{preamble}\n\n{row_text}") > budget:
            flush()
            sections.extend(_split_wide_row(sheet, columns, values, row_number, preamble, budget))
            continue
        candidate = preamble + "\n\n" + "\n\n".join([text for _, text in buffer] + [row_text])
        if buffer and (
            estimate_embedding_tokens(candidate) > budget or len(buffer) >= s.rows_per_chunk_max
        ):
            flush()
        buffer.append((row_number, row_text))
    flush()
    return sections


def _read_xlsx_sheets(path: Path) -> tuple[list[tuple[str, list[list[Any]]]], list[str], dict[str, Any]]:
    """Read visible sheets from an .xlsx into plain row lists.

    Single access point for workbook reading (the documented python-calamine
    swap-point). Hidden and very-hidden sheets are skipped by default — they are
    usually scratch pads or stale copies, and indexing them surprises users.
    """
    from openpyxl import load_workbook

    s = config.spreadsheet
    notes: list[str] = []
    details: dict[str, Any] = {"skipped_sheets": [], "truncated_sheets": []}
    sheets: list[tuple[str, list[list[Any]]]] = []
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        for worksheet in workbook.worksheets:
            if getattr(worksheet, "sheet_state", "visible") != "visible":
                details["skipped_sheets"].append(worksheet.title)
                continue
            rows: list[list[Any]] = []
            truncated = False
            for row in worksheet.iter_rows(values_only=True):
                if len(rows) >= s.max_rows + 1:  # +1 leaves room for the header
                    truncated = True
                    break
                values = list(row[: s.max_cols])
                if len(row) > s.max_cols:
                    truncated = True
                if any(_cell_text(v) for v in values):
                    rows.append(values)
            if truncated:
                details["truncated_sheets"].append(worksheet.title)
            sheets.append((worksheet.title, rows))
    finally:
        workbook.close()
    if details["skipped_sheets"]:
        notes.append("spreadsheet_hidden_sheets_skipped")
    if details["truncated_sheets"]:
        notes.append("spreadsheet_truncated")
    return sheets, notes, details


def _count_uncached_formulas(path: Path) -> int:
    """Count formula cells that carry no cached result.

    `data_only=True` returns cached values, so a workbook saved by a tool that
    never computed them reads as *empty* — the file looks ingested but holds no
    text. Only worth a second pass when the first one actually saw blanks.
    """
    from openpyxl import load_workbook

    s = config.spreadsheet
    total = 0
    workbook = load_workbook(str(path), read_only=True, data_only=False)
    try:
        for worksheet in workbook.worksheets:
            if getattr(worksheet, "sheet_state", "visible") != "visible":
                continue
            for index, row in enumerate(worksheet.iter_rows(values_only=True)):
                if index > s.max_rows:
                    break
                for value in row[: s.max_cols]:
                    if isinstance(value, str) and value.startswith("="):
                        total += 1
    finally:
        workbook.close()
    return total


CSV_ENCODING_SAMPLE_BYTES = 65_536


def _decode_csv_sample(raw: bytes, encoding: str) -> str:
    """Strictly decode a bounded sample without rejecting a split final codepoint."""
    import codecs

    decoder = codecs.getincrementaldecoder(encoding)(errors="strict")
    return decoder.decode(raw, final=False)


def _detect_csv_encoding(raw: bytes) -> tuple[str, dict[str, Any]]:
    """Choose an encoding from a bounded sample and record the decision.

    The UTF-8 result is provisional until the streaming parser reaches EOF. A
    later decode failure triggers one bounded re-detection from the failing
    decoder buffer and a full streaming re-read.
    """
    import codecs

    if raw.startswith(codecs.BOM_UTF8):
        return "utf-8-sig", {"encoding": "utf-8-sig", "source": "bom"}
    for bom, name in ((codecs.BOM_UTF16_LE, "utf-16"), (codecs.BOM_UTF16_BE, "utf-16")):
        if raw.startswith(bom):
            return name, {"encoding": name, "source": "bom"}
    try:
        _decode_csv_sample(raw, "utf-8")
        return "utf-8", {"encoding": "utf-8", "source": "strict"}
    except UnicodeDecodeError:
        pass
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(raw).best()
        if best is not None and best.encoding:
            _decode_csv_sample(raw, best.encoding)
            return best.encoding, {
                "encoding": best.encoding,
                "source": "detected",
                "replacements": 0,
            }
    except Exception:
        logger.exception("csv_encoding_detection_failed")
    return "utf-8", {"encoding": "utf-8", "source": "replace", "replacements": 0}


def _same_encoding(left: str, right: str) -> bool:
    import codecs

    try:
        return codecs.lookup(left).name == codecs.lookup(right).name
    except LookupError:
        return left.casefold() == right.casefold()


def _sniff_csv_delimiter(raw_sample: bytes, encoding: str, errors: str) -> str:
    import csv as csv_module

    try:
        if errors == "strict":
            text = _decode_csv_sample(raw_sample, encoding)
        else:
            text = raw_sample.decode(encoding, errors=errors)
        return csv_module.Sniffer().sniff(text, delimiters=",;\t|").delimiter
    except (UnicodeDecodeError, LookupError, csv_module.Error):
        return ","


def _stream_csv_rows(
    path: Path, encoding: str, errors: str, delimiter: str
) -> tuple[list[list[Any]], bool, int]:
    """Parse CSV incrementally while still consuming the file through EOF.

    Continuing after the row cap is deliberate: a sampled UTF-8 head is only a
    guess until the incremental decoder verifies every later byte. Stored rows
    remain bounded, while a late failure can restart with a detected encoding.
    """
    import csv as csv_module
    import io

    s = config.spreadsheet
    rows: list[list[Any]] = []
    truncated = False
    replacements = 0
    with path.open("rb") as binary_file:
        with io.TextIOWrapper(
            binary_file, encoding=encoding, errors=errors, newline=""
        ) as text_file:
            for values in csv_module.reader(text_file, delimiter=delimiter):
                replacements += sum(value.count("�") for value in values)
                if len(rows) >= s.max_rows + 1:
                    truncated = True
                    # The row cap is already known. Drain decoded text only to
                    # verify the remaining bytes without paying csv parsing cost
                    # for rows that cannot be retained.
                    for remaining in iter(
                        lambda: text_file.read(CSV_ENCODING_SAMPLE_BYTES), ""
                    ):
                        replacements += remaining.count("�")
                    break
                trimmed = values[: s.max_cols]
                if len(values) > s.max_cols:
                    truncated = True
                if any(_cell_text(value) for value in trimmed):
                    rows.append(list(trimmed))
    return rows, truncated, replacements


def _read_csv_sheet(
    path: Path,
) -> tuple[list[tuple[str, list[list[Any]]]], list[str], dict[str, Any]]:
    """Read one CSV sheet with bounded detection and incremental decoding."""
    with path.open("rb") as binary_file:
        raw_sample = binary_file.read(CSV_ENCODING_SAMPLE_BYTES)

    encoding, encoding_details = _detect_csv_encoding(raw_sample)

    def parse(
        selected_encoding: str, selected_details: dict[str, Any]
    ) -> tuple[list[list[Any]], bool, str, dict[str, Any]]:
        errors = "replace" if selected_details["source"] == "replace" else "strict"
        delimiter = _sniff_csv_delimiter(raw_sample, selected_encoding, errors)
        rows, truncated, replacements = _stream_csv_rows(
            path, selected_encoding, errors, delimiter
        )
        if selected_details["source"] in {"detected", "replace"}:
            selected_details = {**selected_details, "replacements": replacements}
        return rows, truncated, delimiter, selected_details

    try:
        rows, truncated, delimiter, encoding_details = parse(encoding, encoding_details)
    except UnicodeDecodeError as exc:
        if encoding_details["source"] == "bom":
            raise
        failure_sample = bytes(exc.object)
        # The traceback retains _stream_csv_rows' partially built rows. Release
        # it before the retry so the failed and successful passes do not overlap
        # in memory.
        exc.__traceback__ = None
        fallback_encoding, fallback_details = _detect_csv_encoding(failure_sample)
        if fallback_details["source"] == "detected" and not _same_encoding(
            fallback_encoding, encoding
        ):
            try:
                rows, truncated, delimiter, encoding_details = parse(
                    fallback_encoding, fallback_details
                )
            except UnicodeDecodeError as fallback_exc:
                fallback_exc.__traceback__ = None
                rows, truncated, delimiter, encoding_details = parse(
                    "utf-8",
                    {"encoding": "utf-8", "source": "replace", "replacements": 0},
                )
        else:
            rows, truncated, delimiter, encoding_details = parse(
                "utf-8",
                {"encoding": "utf-8", "source": "replace", "replacements": 0},
            )

    notes: list[str] = []
    sheet_name = path.stem
    details: dict[str, Any] = {
        "skipped_sheets": [],
        "truncated_sheets": [],
        "encoding": encoding_details,
        "delimiter": delimiter,
    }
    if encoding_details.get("source") in {"detected", "replace"}:
        notes.append("csv_encoding_fallback")
    if truncated:
        details["truncated_sheets"].append(sheet_name)
        notes.append("spreadsheet_truncated")
    return [(sheet_name, rows)], notes, details


def _extract_spreadsheet(path: Path) -> ExtractionResult:
    """Extract an .xlsx / .csv into row-shaped chunks (A6c).

    Returns ``pre_chunked=True``: each section is already the unit we want
    embedded, so `process_source` must not run it through `chunk_sections`.
    """
    _guard_source_size(path)
    s = config.spreadsheet
    suffix = path.suffix.lower()
    if suffix == ".csv":
        sheets, notes, details = _read_csv_sheet(path)
        extractor = "csv"
    else:
        sheets, notes, details = _read_xlsx_sheets(path)
        extractor = "xlsx"

    sections: list[tuple[str, str]] = []
    sheet_reports: list[dict[str, Any]] = []
    saw_blank_cells = False
    for sheet_name, rows in sheets:
        if not rows:
            sheet_reports.append({"sheet": sheet_name, "type": "empty", "rows": 0, "chunks": 0})
            continue
        has_header = _looks_like_header(rows[: max(1, s.header_sample_rows)])
        if has_header:
            columns = [
                _cell_text(v) or _column_label(i) for i, v in enumerate(rows[0])
            ]
            body_rows = [(index + 2, values) for index, values in enumerate(rows[1:])]
        else:
            width = max(len(r) for r in rows)
            columns = [_column_label(i) for i in range(width)]
            body_rows = [(index + 1, values) for index, values in enumerate(rows)]
        saw_blank_cells = saw_blank_cells or any(
            not _cell_text(v) for _, values in body_rows for v in values
        )

        qa = _detect_qa_columns(columns, has_header)
        if qa is not None:
            sheet_sections = _qa_sections(sheet_name, columns, body_rows, qa)
            sheet_type = "qa_pairs"
        else:
            sheet_sections = _record_sections(sheet_name, columns, body_rows)
            sheet_type = "records"
        sections.extend(sheet_sections)
        report = {
            "sheet": sheet_name,
            "type": sheet_type,
            "rows": len(body_rows),
            "columns": len(columns),
            "chunks": len(sheet_sections),
            "header": "detected" if has_header else "generated",
        }
        if qa is not None and qa["auto_detected"]:
            report["qa_columns"] = "auto_detected"
        if len(columns) > s.wide_sheet_cols:
            report["wide"] = True
            if "spreadsheet_wide_sheet" not in notes:
                notes.append("spreadsheet_wide_sheet")
        sheet_reports.append(report)

    if extractor == "xlsx" and saw_blank_cells:
        # Only pay for the second pass when blanks exist — no blanks means no
        # uncached formulas to find.
        try:
            formulas = _count_uncached_formulas(path)
        except Exception:
            logger.exception("formula_scan_failed path=%s", path.name)
            formulas = 0
        if formulas:
            details["uncached_formulas"] = formulas
            notes.append("spreadsheet_uncached_formulas")

    details["sheets"] = sheet_reports
    if not any(report.get("type") == "qa_pairs" for report in sheet_reports):
        # Be explicit that nothing was recognised as Q&A rather than letting the
        # user assume the sheet was understood.
        notes.append("spreadsheet_generic_records")
    logger.info(
        "extract_completed path=%s sections=%s extractor=%s sheets=%s",
        path.name, len(sections), extractor, len(sheet_reports),
    )
    return ExtractionResult(
        sections=sections,
        extractor=extractor,
        notes=notes,
        details=details,
        pre_chunked=True,
    )


def is_mostly_cjk(text: str, threshold: float = 0.30) -> bool:
    """Return True when CJK characters dominate the text (>= threshold).

    CJK and Latin script have very different character density (one CJK char
    carries roughly two English words of meaning), so chunk-size targets and
    sentence splitting both branch on this signal.
    """
    if not text:
        return False
    cjk = len(_CJK_RE.findall(text))
    return cjk / max(1, len(text)) >= threshold


def split_sentences(text: str) -> list[str]:
    """Split text into trimmed sentences, keeping the terminator punctuation."""
    if not text:
        return []
    sentences: list[str] = []
    start = 0
    for match in _SENTENCE_BOUNDARY_RE.finditer(text):
        end = match.end()
        piece = text[start:end].strip()
        if piece:
            sentences.append(piece)
        start = end
    tail = text[start:].strip()
    if tail:
        sentences.append(tail)
    return sentences


def _split_long_sentence(sentence: str, target_chars: int) -> list[str]:
    """Break a sentence that exceeds the chunk target.

    Tries soft punctuation (commas, semicolons) first; if even that leaves
    pieces too large (e.g. a wall of CJK with no internal punctuation), hard
    cuts at target_chars boundaries. Output pieces are <= target_chars.
    """
    pieces: list[str] = []
    buf = ""
    for fragment in _SOFT_BREAK_RE.split(sentence):
        candidate = buf + fragment
        if len(candidate) >= target_chars:
            if buf.strip():
                pieces.append(buf.strip())
            buf = fragment
        else:
            buf = candidate
    if buf.strip():
        pieces.append(buf.strip())

    # Any piece still over budget gets hard-cut. This is the worst case but
    # ensures we never feed a >>target_chars chunk to the embedding API.
    final: list[str] = []
    for piece in pieces:
        while len(piece) > target_chars:
            final.append(piece[:target_chars].strip())
            piece = piece[target_chars:]
        if piece.strip():
            final.append(piece.strip())
    return final


def _section_kind(location: str) -> str:
    """Classify extractor locations so unrelated section kinds do not merge."""
    label = (location or "").lower()
    if " table" in f" {label}" or label.startswith("table"):
        return "table"
    if "header" in label:
        return "header"
    if "footer" in label:
        return "footer"
    if "footnote" in label or "endnote" in label:
        return "footnote"
    if label.endswith(" notes"):
        # A6b: speaker notes are presenter cues, a different register from the
        # slide itself — keep them out of the same chunk so a citation says
        # which one an answer came from.
        return "slide_notes"
    if "text box" in label:
        return "text_box"
    if "transcript" in label:
        return "transcript"
    if label.startswith('sheet "'):  # A6c spreadsheet rows are tabular by nature
        return "table"
    return "body"


def _span_label(locations: list[str]) -> str:
    """Build a citation label for a chunk that may span several sections.

    A chunk packed from one section keeps that section's location verbatim. A
    chunk that merged consecutive sections (e.g. several PDF paragraph blocks
    filled up to the chunk target) is labelled as a first-to-last span so the
    citation still points at the right region. Empty locations are ignored.
    """
    cleaned = [loc for loc in locations if loc]
    if not cleaned:
        return ""
    first, last = cleaned[0], cleaned[-1]
    return first if first == last else f"{first} – {last}"


def chunk_sections(
    sections: list[tuple[str, str]],
    target_chars: int | None = None,
    overlap_sentences: int = DEFAULT_OVERLAP_SENTENCES,
) -> list[tuple[str, str]]:
    """Sentence-aware chunking that packs sentences across sections.

    Same sentence-aware strategy as :func:`chunk_text` (CJK-aware sizing,
    sentence-level overlap, long-sentence fallback), but it fills each chunk up
    to ``target_chars`` worth of sentences **across consecutive sections**
    instead of resetting at every section boundary. Without this, formats whose
    extractor emits many small sections — notably the PDF path's per-paragraph
    blocks (``page N paragraph K``) — leave each short paragraph as its own
    tiny fragment, while single-section formats (TXT/MD) fill to target. Each
    sentence carries its originating ``location`` so the emitted chunk is
    labelled with the source span it covers (see :func:`_span_label`).

    Returns ``(location, chunk_text)`` pairs in document order.
    """
    # Normalise each section's text the same way chunk_text does (collapse the
    # horizontal whitespace PDFs scatter mid-paragraph, keep newlines).
    normalized_sections: list[tuple[str, str]] = []
    for location, text in sections:
        if not text:
            continue
        normalized = _HORIZONTAL_WS_RE.sub(" ", text).strip()
        if normalized:
            normalized_sections.append((location, normalized))
    if not normalized_sections:
        return []

    if target_chars is None:
        combined = "\n".join(text for _, text in normalized_sections)
        target_chars = CJK_TARGET_CHARS if is_mostly_cjk(combined) else LATIN_TARGET_CHARS

    # Flatten to (location, sentence, section_kind) units across all sections.
    units: list[tuple[str, str, str]] = []
    for location, text in normalized_sections:
        kind = _section_kind(location)
        for sentence in split_sentences(text):
            units.append((location, sentence, kind))
    if not units:
        return []

    chunks: list[tuple[str, str]] = []
    current: list[tuple[str, str, str]] = []  # (location, sentence, section_kind)
    current_len = 0

    def flush() -> None:
        if current:
            body = " ".join(sentence for _, sentence, _ in current).strip()
            if body:
                chunks.append((_span_label([loc for loc, _, _ in current]), body))

    for location, sentence, kind in units:
        if current and kind != current[-1][2]:
            flush()
            # Do not carry overlap across body/table/header/footer/etc.; it
            # pollutes citations and can glue unrelated extractor regions.
            current = []
            current_len = 0

        if len(sentence) > target_chars:
            flush()
            # Carry-over does not apply across an over-long sentence — by
            # definition it already contains too much context.
            current = []
            current_len = 0
            for piece in _split_long_sentence(sentence, target_chars):
                if piece:
                    chunks.append((location, piece))
            continue

        if current and current_len + len(sentence) + 1 > target_chars:
            flush()
            if overlap_sentences > 0:
                current = current[-overlap_sentences:]
                current_len = sum(len(s) + 1 for _, s, _ in current)
                # If carrying overlap would make the next chunk exceed the
                # target, drop the overlap. This keeps e5-sized CJK chunks from
                # doubling up around dense boundary sentences.
                if current and current_len + len(sentence) + 1 > target_chars:
                    current = []
                    current_len = 0
            else:
                current = []
                current_len = 0

        current.append((location, sentence, kind))
        current_len += len(sentence) + 1

    flush()
    return [(loc, body) for loc, body in chunks if body]


def chunk_text(
    text: str,
    target_chars: int | None = None,
    overlap_sentences: int = DEFAULT_OVERLAP_SENTENCES,
) -> list[str]:
    """Split a single text into sentence-aware retrieval chunks.

    Strategy:
        1. Normalise horizontal whitespace, keep newlines as boundaries.
        2. Detect CJK-dominance to pick a chunk size (Chinese carries roughly
           2x the information density per character of English, so CJK chunks
           target half the chars).
        3. Split into sentences using ``。！？!?\\n`` as primary boundaries.
        4. Greedily fill chunks up to ``target_chars`` worth of sentences.
        5. Overlap chunks by carrying the last ``overlap_sentences`` sentences
           into the next chunk (sentence-level overlap, not char-level —
           preserves grammar at chunk boundaries).
        6. A single sentence longer than ``target_chars`` is split further by
           soft punctuation, then hard-cut as a last resort.

    Pass ``target_chars=None`` (the default) to auto-pick from text language.
    Thin wrapper over :func:`chunk_sections` for a single unlabelled section.
    """
    return [body for _, body in chunk_sections([("", text)], target_chars, overlap_sentences)]


def estimate_embedding_tokens(text: str) -> int:
    """Rough token count for the embedding input window (A6a).

    An **estimate, not a measurement**: counting real tokens would mean shipping
    the model's tokenizer. CJK-heavy text is charged ~1 token/char (conservative
    — e5 rarely does better) and Latin text ~4 chars/token, reusing the same
    `is_mostly_cjk` split the chunker already relies on. Used only to warn that
    a chunk *may* be truncated, never to change chunking. Ground truth lives in
    `tests/inspect_e5_chunk_tokens.py` (QUALITY.md Q0-5).
    """
    d = config.diagnostics
    per_token = d.cjk_chars_per_token if is_mostly_cjk(text) else d.latin_chars_per_token
    return int(len(text) / per_token) if per_token > 0 else 0


def collect_ingest_diagnostics(
    extraction: "ExtractionResult",
    records: list[tuple[str, str]],
) -> dict[str, Any]:
    """Summarise what an ingest actually produced, for the source preview (A6a).

    Everything here is derived from data the pipeline already had; the point is
    to persist it instead of throwing it away, so a user can tell "extraction
    produced nothing useful" apart from "retrieval/answering went wrong".
    """
    d = config.diagnostics
    section_texts = [text for _, text in extraction.sections]
    chars = sum(len(text) for text in section_texts)
    kinds: dict[str, int] = {}
    for location, _text in extraction.sections:
        kind = _section_kind(location)
        kinds[kind] = kinds.get(kind, 0) + 1
    empty_sections = sum(1 for text in section_texts if not text.strip())
    over_budget = [
        text for _, text in records if estimate_embedding_tokens(text) > d.embedding_token_budget
    ]

    warnings: list[dict[str, Any]] = []
    # Extractor-reported notes become warnings with their own copy; the payload
    # each one needs comes from `details`.
    for note in extraction.notes:
        if note == "pdf_structure_fallback":
            continue  # handled below so it keeps its position in the list
        payload: dict[str, Any] = {"code": note}
        if note == "spreadsheet_hidden_sheets_skipped":
            payload["sheets"] = ", ".join(extraction.details.get("skipped_sheets", []))
        elif note == "spreadsheet_truncated":
            payload["sheets"] = ", ".join(extraction.details.get("truncated_sheets", []))
        elif note == "spreadsheet_uncached_formulas":
            payload["count"] = extraction.details.get("uncached_formulas", 0)
        elif note == "pptx_visual_only_slides":
            payload["count"] = extraction.details.get("slides_without_text", 0)
        elif note == "csv_encoding_fallback":
            encoding = extraction.details.get("encoding", {})
            payload["encoding"] = encoding.get("encoding", "?")
            payload["replacements"] = encoding.get("replacements", 0)
        warnings.append(payload)
    if chars < d.low_text_chars:
        # The scanned-PDF / image-only signal. Deliberately reported even when
        # indexing "succeeded", because that is exactly the silent-failure case.
        warnings.append({"code": "low_text", "chars": chars, "threshold": d.low_text_chars})
    if "pdf_structure_fallback" in extraction.notes:
        warnings.append({"code": "pdf_structure_fallback"})
    if over_budget:
        warnings.append({
            "code": "chunk_over_token_budget",
            "count": len(over_budget),
            "budget": d.embedding_token_budget,
        })
    if empty_sections:
        warnings.append({"code": "empty_sections", "count": empty_sections})

    preview = "\n\n".join(text.strip() for text in section_texts if text.strip())[: d.preview_chars]
    return {
        "extractor": extraction.extractor,
        "chars": chars,
        "sections": len(extraction.sections),
        "chunks": len(records),
        "section_kinds": kinds,
        "warnings": warnings,
        "preview": preview,
        # Format-specific facts (per-sheet type/header decision, CSV encoding).
        # Empty for formats that report none.
        "details": extraction.details,
    }


def get_settings() -> dict[str, Any]:
    """Load the single global LLM settings row with the API key decrypted."""
    with connect() as conn:
        return load_llm_settings(conn) or {}


async def _generate_source_summary(source_id: int) -> None:
    """Generate and persist a per-source TL;DR. Failures are logged only."""
    try:
        with connect() as conn:
            source = conn.execute("SELECT user_id, notebook_id FROM sources WHERE id = ?", (source_id,)).fetchone()
            chunk_rows = conn.execute(
                "SELECT location, text FROM chunks WHERE source_id = ? ORDER BY chunk_index ASC LIMIT 12",
                (source_id,),
            ).fetchall()
            settings = load_llm_settings(conn) or {}
        chunks = [dict(r) for r in chunk_rows]
        if not chunks:
            return
        usage_context = {
            "source_id": source_id,
            "user_id": source["user_id"] if source else None,
            "notebook_id": source["notebook_id"] if source else None,
        }
        summary = await summarize_source(chunks, settings, usage_context=usage_context)
        if not summary:
            return
        with connect() as conn:
            conn.execute(
                "UPDATE sources SET summary = ?, summary_at = CURRENT_TIMESTAMP WHERE id = ?",
                (summary, source_id),
            )
        logger.info("source_summary_persisted source_id=%s chars=%s", source_id, len(summary))
    except Exception:
        logger.exception("source_summary_unhandled source_id=%s", source_id)


async def process_source(source_id: int) -> None:
    """Extract, chunk, embed, and persist vectors for one source record."""
    with connect() as conn:
        source = conn.execute("SELECT * FROM sources WHERE id = ?", (source_id,)).fetchone()
        if source is None:
            logger.warning("ingest_source_missing source_id=%s", source_id)
            return
        # Diagnostics describe the *current* ingest only, so clear the previous
        # run's alongside its chunks — a reindex must not show stale warnings.
        conn.execute(
            "UPDATE sources SET status = 'processing', error = '', diagnostics_json = '{}', "
            "updated_at = CURRENT_TIMESTAMP WHERE id = ?",
            (source_id,),
        )
        conn.execute("DELETE FROM chunks WHERE source_id = ?", (source_id,))
    try:
        delete_source_vectors(source_id, source["user_id"])
    except Exception:
        logger.exception("vector_source_delete_failed source_id=%s", source_id)

    # Which phase we are in, so a failure can say whether extraction produced
    # nothing or the embedding endpoint was down (A6a "failure reason").
    stage = "extract"
    diagnostics: dict[str, Any] = {}
    try:
        logger.info(
            "ingest_started source_id=%s user_id=%s filename=%s",
            source_id,
            source["user_id"],
            source["filename"],
        )
        # Off the event loop: extraction is synchronous, CPU-bound, and can run
        # for tens of seconds (measured: 32s for an 881-section PDF through
        # pdfplumber). With the inline worker that time is spent inside the web
        # process, so running it here would freeze every request — including the
        # 2s source-row polls that are supposed to show "processing", which is
        # exactly how this surfaced: the row sat at "uploaded" because the
        # server could not answer a single poll until extraction finished.
        # `extract_sections` only reads the file and returns text, so it holds
        # no connection and is safe to hand to a thread.
        extraction = await asyncio.to_thread(extract_sections, Path(source["stored_path"]))
        sections = extraction.sections
        stage = "chunk"
        # Pack sentences across sections up to the chunk target so formats that
        # emit many small sections (PDF per-paragraph blocks) produce the same
        # well-sized chunks as single-section formats (TXT/MD) rather than
        # hundreds of tiny fragments. Spreadsheets opt out: a row is already the
        # semantic unit and re-packing would glue unrelated records together.
        records: list[tuple[str, str]] = (
            list(sections) if extraction.pre_chunked else chunk_sections(sections)
        )
        # Collected before the empty-records check so a source that extracted
        # nothing still gets diagnostics — that is the case users most need to
        # see, and it is the only window into a failed source.
        diagnostics = collect_ingest_diagnostics(extraction, records)
        if not records:
            raise ValueError("No extractable text found.")

        logger.info(
            "ingest_chunked source_id=%s sections=%s chunks=%s extractor=%s warnings=%s",
            source_id,
            len(sections),
            len(records),
            diagnostics["extractor"],
            len(diagnostics["warnings"]),
        )
        stage = "embed"
        embeddings = await embed_texts(
            [text for _, text in records],
            get_settings(),
            role="passage",
            usage_context={"user_id": source["user_id"], "notebook_id": source["notebook_id"], "source_id": source_id},
        )
        stage = "store"
        with connect() as conn:
            chunk_rows = [
                (source["user_id"], source_id, index, location, text, dumps(embedding))
                for index, ((location, text), embedding) in enumerate(zip(records, embeddings))
            ]
            conn.executemany(
                """
                INSERT INTO chunks (user_id, source_id, chunk_index, location, text, embedding_json)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                chunk_rows,
            )
            inserted = conn.execute(
                """
                SELECT chunks.*, sources.filename
                FROM chunks JOIN sources ON sources.id = chunks.source_id
                WHERE chunks.source_id = ?
                ORDER BY chunks.chunk_index
                """,
                (source_id,),
            ).fetchall()
            conn.execute(
                "UPDATE sources SET status = 'indexed', error = '', diagnostics_json = ?, "
                "updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                (dumps(diagnostics), source_id),
            )
        upsert_chunks(
            [
                {
                    "id": row["id"],
                    "user_id": row["user_id"],
                    "source_id": row["source_id"],
                    "chunk_index": row["chunk_index"],
                    "filename": row["filename"],
                    "location": row["location"],
                    "text": row["text"],
                    "embedding": embeddings[row["chunk_index"]],
                }
                for row in inserted
            ]
        )
        logger.info("ingest_completed source_id=%s chunks=%s", source_id, len(records))
        # Best-effort per-source summary. Runs AFTER status='indexed' so a
        # summarization failure leaves the source fully usable for retrieval.
        await _generate_source_summary(source_id)
    except Exception as exc:
        # Keep whatever was learned before the failure: which stage broke, and
        # (when extraction got that far) what it managed to pull out. Without
        # this a failed source shows an error string and nothing else.
        failed_diagnostics = {**diagnostics, "failed_stage": stage}
        with connect() as conn:
            conn.execute(
                """
                UPDATE sources
                SET status = 'failed', error = ?, diagnostics_json = ?, updated_at = CURRENT_TIMESTAMP
                WHERE id = ?
                """,
                (str(exc)[:500], dumps(failed_diagnostics), source_id),
            )
        logger.exception("ingest_failed source_id=%s stage=%s", source_id, stage)
