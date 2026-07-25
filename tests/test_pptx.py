"""A6b Phase 1 — text-first PPTX ingestion.

Fixtures are built with python-pptx so the suite carries no binary assets. The
recurring bug class this guards against is *silent content loss*: grouped
shapes, tables, and speaker notes are all easy to walk past.
"""
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.ingest import extract_sections


def _deck(tmp_path: Path) -> Presentation:
    return Presentation()


def _save(prs: Presentation, tmp_path: Path) -> Path:
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def test_slides_become_slide_scoped_sections_with_titles_and_body(tmp_path):
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])   # title + content
    slide.shapes.title.text = "系統架構"
    slide.placeholders[1].text_frame.text = "FastAPI 與 Jinja2\nChroma 向量檢索"
    second = prs.slides.add_slide(prs.slide_layouts[1])
    second.shapes.title.text = "部署考量"
    second.placeholders[1].text_frame.text = "單機部署即可滿足需求"

    result = extract_sections(_save(prs, tmp_path))

    assert result.extractor == "pptx"
    assert result.pre_chunked is False          # slides flow through the chunker
    locations = [location for location, _ in result.sections]
    assert locations == ["slide 1", "slide 2"]
    assert result.sections[0][1].startswith("系統架構")   # title leads the slide text
    assert "Chroma 向量檢索" in result.sections[0][1]
    assert result.details["slides"] == 2


def test_speaker_notes_become_their_own_section(tmp_path):
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "定價"
    slide.notes_slide.notes_text_frame.text = "這裡要提醒客戶合約是年約。"

    result = extract_sections(_save(prs, tmp_path))

    by_location = dict(result.sections)
    assert "slide 1 notes" in by_location
    assert by_location["slide 1 notes"] == "這裡要提醒客戶合約是年約。"
    assert result.details["notes"] == 1


def test_tables_are_extracted_as_their_own_slide_scoped_sections(tmp_path):
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # title only
    slide.shapes.title.text = "比較"
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table = shape.table
    table.cell(0, 0).text = "方案"
    table.cell(0, 1).text = "價格"
    table.cell(1, 0).text = "標準版"
    table.cell(1, 1).text = "1000"

    result = extract_sections(_save(prs, tmp_path))

    by_location = dict(result.sections)
    assert "slide 1 table 1" in by_location
    assert "方案 | 價格" in by_location["slide 1 table 1"]
    assert "標準版 | 1000" in by_location["slide 1 table 1"]
    assert result.details["tables"] == 1


def test_grouped_shapes_are_not_silently_dropped(tmp_path):
    """Grouped text boxes are the PPTX version of the nested-table bug."""
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    box_a = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box_a.text_frame.text = "群組內的第一段文字"
    box_b = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
    box_b.text_frame.text = "群組內的第二段文字"
    # python-pptx has no public grouping API; move the shape elements into a
    # group element the same way PowerPoint stores them.
    group = slide.shapes.add_group_shape([box_a, box_b])
    assert group is not None

    result = extract_sections(_save(prs, tmp_path))

    body = dict(result.sections)["slide 1"]
    assert "群組內的第一段文字" in body
    assert "群組內的第二段文字" in body


def test_visual_only_slides_are_reported_not_silently_empty(tmp_path):
    """A slide of pictures must say so — otherwise the user blames retrieval."""
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image = tmp_path / "dot.png"
    # 1x1 transparent PNG.
    image.write_bytes(bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000a49444154789c6360000002000100ffff03000006000557bfabd400"
        "00000049454e44ae426082"
    ))
    slide.shapes.add_picture(str(image), Inches(1), Inches(1))

    result = extract_sections(_save(prs, tmp_path))

    assert result.sections == []
    assert result.details["images"] == 1
    assert result.details["slides_without_text"] == 1
    assert "pptx_visual_only_slides" in result.notes


def test_speaker_notes_do_not_merge_into_slide_body_chunks(tmp_path):
    """Notes are a different register — a citation must say which one it quoted."""
    from app.ingest import _section_kind, chunk_sections

    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "定價"
    slide.placeholders[1].text_frame.text = "標準版每年新台幣一萬元。"
    slide.notes_slide.notes_text_frame.text = "口頭補充：這個價格還有議價空間，不要寫進簡報。"

    result = extract_sections(_save(prs, tmp_path))
    chunks = chunk_sections(result.sections)

    assert _section_kind("slide 1 notes") == "slide_notes"
    assert _section_kind("slide 1") == "body"
    bodies = [text for _, text in chunks]
    assert not any("議價空間" in text and "標準版每年" in text for text in bodies)


def test_slide_sections_read_body_then_tables(tmp_path):
    """PPTX shape order isn't visual order; emit a predictable per-slide order."""
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "格式支援"
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    shape.table.cell(0, 0).text = "格式"
    shape.table.cell(0, 1).text = "狀態"
    shape.table.cell(1, 0).text = "PPTX"
    shape.table.cell(1, 1).text = "本次新增"

    result = extract_sections(_save(prs, tmp_path))

    assert [location for location, _ in result.sections] == ["slide 1", "slide 1 table 1"]


def test_a_slide_with_only_a_table_is_not_counted_as_visual_only(tmp_path):
    """A table is readable content — it must not trip the 'no text' warning."""
    prs = _deck(tmp_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 1, Inches(1), Inches(1), Inches(3), Inches(1))
    shape.table.cell(0, 0).text = "項目"
    shape.table.cell(1, 0).text = "內容"

    result = extract_sections(_save(prs, tmp_path))

    assert result.details["slides_without_text"] == 0
    assert "pptx_visual_only_slides" not in result.notes
